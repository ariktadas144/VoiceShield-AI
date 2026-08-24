import os, sys
from pptx import Presentation
from PIL import ImageFont
TTF=os.path.expanduser("~/.local/share/fonts/ArchitectsDaughter-Regular.ttf")
PADX,PADY,LINE,GAP=0.11,0.07,1.24,2.0
_fc={}
def _f(px):
    if px not in _fc:_fc[px]=ImageFont.truetype(TTF,px)
    return _fc[px]
def lines(t,sz,bold,w):
    if not t.strip():return 1
    f=_f(max(6,int(round(sz*96/72))));mp=max(8.0,w*96);sw=1.06 if bold else 1.0
    n,cur=1,""
    for wd in t.split():
        st=(cur+" "+wd).strip()
        if f.getlength(st)*sw<=mp:cur=st
        else:n+=1;cur=wd
    return n
prs=Presentation("docs/streaming_bridge/VoiceShield-Bridge-3slides.pptx")
SW,SH=prs.slide_width,prs.slide_height
def I(e):return e/914400.0
def R(s):return (s.left,s.top,s.left+s.width,s.top+s.height)
def ov(a,b):
    x=max(0,min(a[2],b[2])-max(a[0],b[0]));y=max(0,min(a[3],b[3])-max(a[1],b[1]));return x*y
bad=0
for n,sl in enumerate(prs.slides,1):
    shapes=list(sl.shapes)
    conns=[x for x in shapes if type(x).__name__=="Connector"]
    boxes=[x for x in shapes if x.shape_type==1]
    for sh in shapes:
        l,t,r,b=R(sh)
        if l<-2000 or t<-2000 or r>SW+2000 or b>SH+2000:
            print(f"  s{n} OFF-SLIDE {I(l):.2f},{I(t):.2f}->{I(r):.2f},{I(b):.2f} '{sh.text_frame.text[:22] if sh.has_text_frame else ''}'");bad+=1
        if sh.has_text_frame and sh.text_frame.text.strip():
            isbox=sh.shape_type==1
            iw=I(sh.width)-(2*PADX if isbox else 0)
            need=(2*PADY if isbox else 0)
            for p in sh.text_frame.paragraphs:
                rs=list(p.runs)
                if not rs: need+=8*LINE/72; continue
                txt="".join(r.text for r in rs)
                sz=max((r.font.size.pt if r.font.size else 11) for r in rs)
                bo=any(r.font.bold for r in rs)
                need+=lines(txt,sz,bo,iw)*sz*LINE/72+GAP/72
            if need>I(sh.height)+0.015:
                print(f"  s{n} TEXT OVERFLOW +{need-I(sh.height):.3f}in '{sh.text_frame.text[:36]}'");bad+=1
    for i,a in enumerate(boxes):
        for b2 in boxes[i+1:]:
            if ov(R(a),R(b2))>0:
                print(f"  s{n} BOX OVERLAP '{a.text_frame.text[:16]}'/'{b2.text_frame.text[:16]}'");bad+=1
    labels=[x for x in shapes if x.shape_type==17 and 0<len(x.text_frame.text.strip())<14]
    for lb in labels:
        for bx in boxes:
            if ov(R(lb),R(bx))>0:
                print(f"  s{n} LABEL '{lb.text_frame.text.strip()}' overlaps a box");bad+=1
    ln=[(I(x.width)**2+I(x.height)**2)**0.5 for x in conns]
    print(f"slide{n}: {len(shapes)} shapes · {len(conns)} arrows"+(f" (shortest {min(ln):.2f}in)" if ln else "")+f" · bottom {max(I(R(x)[3]) for x in shapes):.2f}/7.50")
print("\nISSUES:",bad,"->","PASS" if bad==0 else "FAIL")
sys.exit(0 if bad==0 else 1)
