"""Build the 3 review slides for the Streaming Middleware Bridge.

Every element is a NATIVE PowerPoint shape — rounded rectangles, connectors and
text runs. No raster images. Click any box, arrow or word to edit it.

Two things make this reliable rather than hopeful:

1. Text is MEASURED, not estimated. Box heights come from PIL wrapping the real
   font at the real point size, so nothing overflows its card. Architects
   Daughter runs ~6% wider than Arial, which is exactly the kind of thing that
   silently breaks a layout when you guess.
2. The hand-drawn outline is PowerPoint's own sketch line style (the 2018
   "sketchyshapes" extension). Consumers that don't understand it render a
   clean line, so it degrades safely.

FONT: Architects Daughter. Install fonts/ArchitectsDaughter-Regular.ttf on any
machine that opens this deck, or PowerPoint will substitute and reflow.

Re-run after editing:  .venv/bin/python docs/streaming_bridge/build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from PIL import ImageFont

# ── palette (Excalidraw defaults) ──────────────────────────────────────────
INK  = RGBColor(0x1E,0x1E,0x1E); MUT = RGBColor(0x5C,0x56,0x50)
BLUE = RGBColor(0x19,0x71,0xC2); BLUEF = RGBColor(0xA5,0xD8,0xFF); BLUEP = RGBColor(0xE7,0xF2,0xFC)
RED  = RGBColor(0xE0,0x31,0x31); REDF  = RGBColor(0xFF,0xC9,0xC9); REDP  = RGBColor(0xFF,0xF0,0xEF)
GRN  = RGBColor(0x2F,0x9E,0x44); GRNF  = RGBColor(0xB2,0xF2,0xBB); GRNP  = RGBColor(0xEE,0xFB,0xF1)
AMB  = RGBColor(0xF0,0x8C,0x00); AMBF  = RGBColor(0xFF,0xEC,0x99); AMBP  = RGBColor(0xFF,0xF8,0xE6)
VIO  = RGBColor(0x9C,0x36,0xB5); VIOF  = RGBColor(0xEE,0xBE,0xFA); VIOP  = RGBColor(0xF8,0xEF,0xFB)
PANEL= RGBColor(0xF5,0xF2,0xEC); WHITE= RGBColor(0xFF,0xFF,0xFF)

F   = "Architects Daughter"
TTF = os.path.expanduser("~/.local/share/fonts/ArchitectsDaughter-Regular.ttf")
PADX, PADY = 0.11, 0.07          # text-frame margins, inches
LINE, GAP  = 1.24, 2.0           # line-height multiple, per-paragraph gap in pt
SKNS = "http://schemas.microsoft.com/office/drawing/2018/sketchyshapes"
_seed = [1799159648]
_fc = {}

def _font(px):
    if px not in _fc: _fc[px] = ImageFont.truetype(TTF, px)
    return _fc[px]

def _lines(text, size_pt, bold, width_in):
    """How many wrapped lines this text takes at this size in this width."""
    if not text.strip(): return 1
    f = _font(max(6, int(round(size_pt * 96 / 72))))
    maxpx = max(8.0, width_in * 96)
    swell = 1.06 if bold else 1.0     # PowerPoint fakes bold on a single-weight face
    n, cur = 1, ""
    for w in text.split():
        t = (cur + " " + w).strip()
        if f.getlength(t) * swell <= maxpx: cur = t
        else: n += 1; cur = w
    return n

def blocks_h(blocks, width_in):
    """Height in inches needed to render these paragraphs in this width."""
    inner = width_in - 2 * PADX
    h = 2 * PADY
    for runs in blocks:
        if not runs: h += 8 * LINE / 72; continue
        txt  = "".join(r[0] for r in runs)
        size = max(r[1] for r in runs)
        bold = any(r[2] for r in runs)
        h += _lines(txt, size, bold, inner) * size * LINE / 72 + GAP / 72
    return h

def sketchy(shape, kind="lineSketchFreehand"):
    ln = shape.line._get_or_add_ln()
    ext_lst = ln.find(qn('a:extLst'))
    if ext_lst is None: ext_lst = etree.SubElement(ln, qn('a:extLst'))
    ext = etree.SubElement(ext_lst, qn('a:ext'))
    ext.set('uri', '{C807C97D-BFC1-408E-A445-0C87EB9F89A2}')
    p = etree.SubElement(ext, '{%s}lineSketchStyleProps' % SKNS)
    _seed[0] += 7919; p.set('sd', str(_seed[0]))
    t = etree.SubElement(p, '{%s}type' % SKNS)
    etree.SubElement(t, '{%s}%s' % (SKNS, kind))
    return shape

def box(s, x, y, w, h, line=INK, fill=WHITE, lw=2.0, radius=0.05,
        anchor=MSO_ANCHOR.MIDDLE, sketch=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    try: sh.adjustments[0] = radius
    except Exception: pass
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(PADX)
    tf.margin_top = tf.margin_bottom = Inches(PADY)
    tf.vertical_anchor = anchor
    if sketch: sketchy(sh)
    return sh

def rich(shape, blocks, align=None):
    tf = shape.text_frame; tf.paragraphs[0].text = ""
    for i, runs in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i: p.space_before = Pt(GAP)
        if align is not None: p.alignment = align
        for (txt, size, bold, color) in runs:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = F
    return shape

def autobox(s, x, y, w, blocks, line=INK, fill=WHITE, lw=2.0, min_h=0.0,
            anchor=MSO_ANCHOR.TOP, align=None):
    """Box sized to exactly fit its measured content."""
    h = max(min_h, blocks_h(blocks, w))
    sh = box(s, x, y, w, h, line, fill, lw, anchor=anchor)
    rich(sh, blocks, align)
    return sh, h

def tbox(s, x, y, w, h):
    b = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return b

def free(s, x, y, w, blocks, align=None):
    h = blocks_h(blocks, w + 2 * PADX)
    b = tbox(s, x, y, w, h)
    rich(b, blocks, align)
    return b, h

def arrow(s, x1, y1, x2, y2, color=INK, lw=3.0):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(lw)
    ln = c.line._get_or_add_ln()
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'triangle'); te.set('w', 'lg'); te.set('len', 'lg')
    return c

def caption(s, cx, y, txt, color=MUT, size=8, w=0.64):
    """Centred label sitting in clear space, never on top of a shape.

    Width is capped to the gap it occupies — a label wider than its gap lands
    on the neighbouring box and the text becomes unreadable.
    """
    x = min(max(cx - w / 2, 0.06), 13.27 - w)
    b = tbox(s, x, y, w, 0.20)
    p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = F
    return b

def heading(s, title, tail, sub):
    hb = [[(title, 24, True, INK), (tail, 24, False, BLUE)]]
    b = tbox(s, 0.42, 0.14, 12.5, blocks_h(hb, 12.5 + 2 * PADX))
    rich(b, hb)
    sb = [[(sub, 10.5, False, MUT)]]
    sh = blocks_h(sb, 12.5 + 2 * PADX)
    b2 = tbox(s, 0.42, 0.14 + blocks_h(hb, 12.5 + 2 * PADX) + 0.04, 12.5, sh)
    rich(b2, sb)
    return 0.14 + blocks_h(hb, 12.5 + 2 * PADX) + 0.04 + sh

prs = Presentation(); prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]
print("helpers ready · font:", F)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Role & interfaces
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
HB1 = heading(s1, "Streaming Middleware Bridge", "  —  Component 2 of 4",
  "The only stateful, real-time component. Turns live forked call audio into scored analysis windows inside a 500 ms budget, "
  "for ~1,200 concurrent calls — on a stream that cannot be paused.")

# wide gaps so the connectors are long enough to read and label
CY = HB1 + 0.14
ctx = [
 (0.42, 2.45, INK, PANEL, 2.25,
  [[("(1) TELEPHONY &", 13, True, INK)], [("MEDIA GATEWAY", 13, True, INK)],
   [("FreeSWITCH + mod_audio_stream", 8.5, False, MUT)],
   [("SIP / RTP · forks call audio", 8.5, False, MUT)],
   [("thread-isolated per channel", 8.5, False, MUT)]]),
 (3.55, 3.10, BLUE, BLUEF, 3.5,
  [[("(2) STREAMING BRIDGE", 15, True, BLUE)], [("gRPC", 15, True, BLUE)],
   [("MY COMPONENT", 9.5, True, BLUE)],
   [("normalise · session state", 8.5, False, INK)],
   [("buffer · window · transport", 8.5, False, INK)],
   [("reorder · fan-out · degrade", 8.5, False, INK)]]),
 (7.33, 2.55, INK, PANEL, 2.25,
  [[("(3) AI & DSP ENGINE", 13, True, INK)],
   [("WavLM + LFCC/MGD fusion", 8.5, False, MUT)],
   [("Silero VAD · ECAPA-TDNN", 8.5, False, MUT)],
   [("RTF <= 0.1 target", 8.5, False, MUT)],
   [("dev EER 0.55% in-domain", 8.5, False, MUT)]]),
 (10.56, 2.35, INK, PANEL, 2.25,
  [[("(4) RISK ENGINE", 13, True, INK)], [("+ DASHBOARD", 13, True, INK)],
   [("weighted multi-signal fusion", 8.5, False, MUT)],
   [("dynamic thresholds", 8.5, False, MUT)],
   [("prevention + CRM webhooks", 8.5, False, MUT)]]),
]
CTX_H = max(blocks_h(b, w) for _, w, _, _, _, b in ctx) + 0.95
for x, w, ln_c, fl_c, lw, blk in ctx:
    yy = CY - 0.10 if ln_c == BLUE else CY
    hh = CTX_H + 0.20 if ln_c == BLUE else CTX_H
    sh = box(s1, x, yy, w, hh, ln_c, fl_c, lw, anchor=MSO_ANCHOR.MIDDLE)
    rich(sh, blk, PP_ALIGN.CENTER)
MID = CY + CTX_H / 2

arrow(s1, 2.95, MID, 3.47, MID, INK, 3.25)
caption(s1, 3.21, MID - 0.36, "L16", MUT, 8, 0.56)
arrow(s1, 6.73, MID - 0.26, 7.25, MID - 0.26, INK, 3.25)
caption(s1, 6.99, MID - 0.62, "windows", MUT, 8, 0.56)
arrow(s1, 7.25, MID + 0.26, 6.73, MID + 0.26, VIO, 3.25)
caption(s1, 6.99, MID + 0.34, "scores", VIO, 8, 0.56)
arrow(s1, 9.96, MID, 10.48, MID, INK, 3.25)
caption(s1, 10.22, MID - 0.36, "results", MUT, 8, 0.56)

CB = CY + CTX_H + 0.22
free(s1, 0.42, CB, 12.5, [[("FOUR INTERFACE CONTRACTS", 10, True, BLUE),
  ("   — agree these before anyone writes code; changing one later costs all four of us", 9.5, False, MUT)]])

contracts = [
 (BLUE, BLUEP, "CONTRACT A", "Gateway  ->  Bridge",
  [("session_id","server-issued or HMAC-signed"),("encoding","L16_LE / L16_BE / PCMU / PCMA"),
   ("sample_rate","8000 or 16000 — never assumed"),("channels + leg","mono only; 2 legs = 2 sessions"),
   ("sample_offset","the authoritative clock"),("concealed","true if gateway invented the audio"),
   ("consent","DPDP gate — no consent, no audio"),
   ("OPEN Q","is your L16 big- or little-endian? Byte-swapped PCM is loud noise and the detector will confidently call it synthetic.")]),
 (VIO, VIOP, "CONTRACT B", "Bridge  <->  AI Engine",
  [("AudioChunk","pcm, window_seq, sample_offset, gap_before, partial, session_epoch"),
   ("RiskScoreUpdate","window_seq ECHOED back, model_version, status"),
   ("why echoed","out-of-order results are otherwise unrecoverable on my side"),
   ("batch","always 1 — Dhwani silently collapses batches to one result"),
   ("deadline","400 ms; on breach the window is abandoned, never retried"),
   ("pooling","80 streams/conn — HTTP/2 caps at 100")]),
 (GRN, GRNP, "CONTRACT C", "Bridge  ->  Risk Engine",
  [("ordering","monotonic window_seq, deduplicated"),
   ("gaps","gap_before=true resets any smoothing"),
   ("status","OK / SKIPPED_SILENCE / DEGRADED / DETECTOR_UNAVAILABLE"),
   ("critical","a missing signal must renormalise, never be substituted with 0.0"),
   ("OPEN Q","on DETECTOR_UNAVAILABLE do you renormalise or hold the last score?"),
   ("cadence","10 scores/sec at L0, fewer when degraded")]),
 (AMB, AMBP, "CONTRACT D", "Bridge  ->  Dashboard",
  [("transport","best-effort WebSocket"),("delivery","lossy by design"),
   ("rule","never in the critical path; must never apply backpressure that reaches audio"),
   ("why lossy","a dashboard that can slow the audio path is a dashboard that can drop fraud detection"),
   ("reconnect","client resumes, no server state"),
   ("cardinality","no per-session metric labels")]),
]
blks = []
for ln_c, fl_c, cid, who, rows in contracts:
    b = [[(cid, 9.5, True, ln_c)], [(who, 10.5, True, INK)]]
    for k, v in rows:
        b.append([(k + "  ", 8, True, INK), (v, 8, False, MUT)])
    blks.append(b)
CW = 2.98
CH2 = max(blocks_h(b, CW) for b in blks)
x = 0.42
for (ln_c, fl_c, *_), b in zip(contracts, blks):
    sh = box(s1, x, CB + 0.30, CW, CH2, ln_c, fl_c, 2.0, anchor=MSO_ANCHOR.TOP)
    rich(sh, b); x += 3.14

CC = CB + 0.30 + CH2 + 0.16
cb = [[("CORE CONSTRAINT   ", 10, True, RED),
       ("Every other service in this system can tell its upstream to slow down. A phone call cannot be paused.", 10, False, INK)],
      [("Everything else — the ring buffer, the drop-oldest policy, the adaptive hop, the whole degradation ladder — is a consequence of that one fact.", 9, False, MUT)]]
box(s1, 0.42, CC, 12.49, blocks_h(cb, 12.49), RED, REDP, 2.0, anchor=MSO_ANCHOR.MIDDLE)
rich(s1.shapes[-1], cb)
print("slide 1 ok  bottom=%.2f" % (CC + blocks_h(cb, 12.49)))

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Internal pipeline
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
HB2 = heading(s2, "Inside the Bridge", "  —  Eight Stages, One Timeline",
  "N dissimilar transports become one canonical PCM timeline; that timeline becomes fixed windows the model was trained to expect.")

ing = [("FreeSWITCH", "WebSocket · L16 binary", "8 or 16 kHz · primary path"),
       ("Asterisk", "TCP · AudioSocket", "length-prefixed framing"),
       ("Twilio", "WSS · base64 JSON", "G.711 mu-law 8 kHz")]
ING_W = 1.95
ing_blk = [[[(n, 11, True, INK)], [(a, 8, False, MUT)], [(b, 8, False, MUT)]] for n, a, b in ing]
ING_H = max(blocks_h(b, ING_W) for b in ing_blk) + 0.52
iy = HB2 + 0.14
icenters = []
for b in ing_blk:
    sh = box(s2, 0.42, iy, ING_W, ING_H, INK, WHITE, 2.0, anchor=MSO_ANCHOR.MIDDLE)
    rich(sh, b, PP_ALIGN.CENTER); icenters.append(iy + ING_H / 2); iy += ING_H + 0.20

stages = [
 ("INGRESS ADAPTERS", BLUE, BLUEF, [("N formats -> 1 frame", True), ("decode · mono · 16 kHz", False),
   ("soxr VHQ resample", False), ("endianness validated", False), ("same path as training", False)]),
 ("SESSION MANAGER", BLUE, BLUEF, [("lifecycle + fencing", True), ("Redis epoch token", False),
   ("consent gate (DPDP)", False), ("resume within 15 s", False), ("2nd claimant rejected", False)]),
 ("RING BUFFER", BLUE, BLUEF, [("bounded 256 KB", True), ("jitter absorb 60 ms", False),
   ("drop OLDEST on full", False), ("never blocks ingest", False), ("1 producer, 1 consumer", False)]),
 ("WINDOWER", BLUE, BLUEF, [("2.0 s / 100 ms hop", True), ("sample-accurate", False),
   ("VAD-gated (~40% saved)", False), ("length from checkpoint", False), ("hop by samples", False)]),
 ("gRPC POOL", VIO, VIOF, [("bidi stream · batch=1", True), ("80 streams per conn", False),
   ("400 ms deadline", False), ("retry on NEW stream", False), ("keepalive < LB idle", False)]),
]
SW_, SGAP, SX0 = 1.57, 0.52, 2.95
st_blk = []
for nm, ln_c, fl_c, rows in stages:
    b = [[(nm, 10.5, True, ln_c)]]
    for txt, bold in rows: b.append([(txt, 8, bold, INK if bold else MUT)])
    st_blk.append(b)
ST_H = max(blocks_h(b, SW_) for b in st_blk) + 0.90
SY = HB2 + 0.25
sx = SX0
for (nm, ln_c, fl_c, rows), b in zip(stages, st_blk):
    sh = box(s2, sx, SY, SW_, ST_H, ln_c, fl_c, 3.0, anchor=MSO_ANCHOR.TOP)
    rich(sh, b, PP_ALIGN.CENTER)
    if sx + SW_ + SGAP < 12.9:
        arrow(s2, sx + SW_ + 0.05, SY + ST_H / 2, sx + SW_ + SGAP - 0.05, SY + ST_H / 2, INK, 3.0)
    sx += SW_ + SGAP
for c in icenters:
    arrow(s2, 2.43, c, 2.89, SY + ST_H / 2, INK, 2.5)

RY = SY + ST_H + 0.46
ret = [("REORDER", "by window_seq", "300 ms bounded wait", BLUE, BLUEF),
       ("FAN-OUT", "dedupe · annotate gaps", "risk / dash / CRM", BLUE, BLUEF),
       ("(4) RISK ENGINE", "weighted fusion", "prevention actions", INK, PANEL)]
ret_blk = [[[(n, 10, True, c if c != INK else INK)], [(a, 8, False, MUT)], [(b, 8, False, MUT)]]
           for n, a, b, c, _ in ret]
RET_H = max(blocks_h(b, SW_) for b in ret_blk) + 0.18
rx = SX0 + 4 * (SW_ + SGAP)
last_stage_cx = rx + SW_ / 2
for (n, a, b, ln_c, fl_c), blk in zip(ret, ret_blk):
    sh = box(s2, rx, RY, SW_, RET_H, ln_c, fl_c, 2.5, anchor=MSO_ANCHOR.MIDDLE)
    rich(sh, blk, PP_ALIGN.CENTER)
    if rx - SW_ - SGAP > 2.0:
        arrow(s2, rx - 0.05, RY + RET_H / 2, rx - SGAP + 0.05, RY + RET_H / 2, VIO, 3.0)
    rx -= SW_ + SGAP
arrow(s2, last_stage_cx, SY + ST_H + 0.06, last_stage_cx, RY - 0.06, VIO, 3.0)
caption(s2, last_stage_cx - 0.62, SY + ST_H + 0.14, "scores", VIO, 8, 0.56)

DY = RY + RET_H + 0.34
free(s2, 0.42, DY, 12.5, [[("THREE DECISIONS THAT DETERMINE WHETHER THIS SCALES", 10, True, BLUE)]])
dec = [
 (BLUE, BLUEP, "CONCURRENCY", "One asyncio loop. No per-session threads.",
  "1,200 threads costs ~9.6 GB of stack before any audio is buffered, and collapses on context switching. Sessions are coroutines; each ring has exactly one producer and one consumer, so it needs no lock.",
  "Verified by: 1,200 concurrent sessions held for 30 min inside the latency budget."),
 (VIO, VIOP, "THE CLOCK", "sample_offset is authoritative — not wall time.",
  "Wall clocks drift, jump on NTP and differ per pod. A gateway running at 8000.5 Hz against an assumed 8000 Hz accumulates 1.8 s of error per hour. Windowing must be reproducible.",
  "Verified by: sample_offset stays monotonic across an injected 300 ms gap and a reconnect."),
 (GRN, GRNP, "MEMORY", "256 KB per session, enforced.",
  "4 s ring x 16 kHz x 4 bytes = 256 KB. At 1,200 sessions that is ~300 MB. The same design with an unbounded per-session accumulator over a 1-hour call would be 690 GB.",
  "Verified by: 8-hour soak at 200 sessions — any upward memory slope is a per-session leak."),
]
dec_blk = [[[(t, 9, True, c)], [(hl, 10, True, INK)], [(bd, 8, False, INK)], [(vf, 8, True, c)]]
           for c, _, t, hl, bd, vf in dec]
DW = 4.06
DH = max(blocks_h(b, DW) for b in dec_blk)
dx = 0.42
for (ln_c, fl_c, *_), b in zip(dec, dec_blk):
    sh = box(s2, dx, DY + 0.28, DW, DH, ln_c, fl_c, 2.0, anchor=MSO_ANCHOR.TOP)
    rich(sh, b); dx += 4.22
print("slide 2 ok  bottom=%.2f" % (DY + 0.28 + DH))

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Latency, degradation, silent failures
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
HB3 = heading(s3, "Latency, Failure Handling", "  &  the Edge Cases That Fail Silently",
  "Most failures are loud, and therefore cheap. The expensive ones produce plausible-looking numbers that are wrong.")

LW_ = 6.16
free(s3, 0.42, HB3 + 0.14, LW_, [[("500 ms END-TO-END SLA", 9.5, True, INK),
                            ("   caller speaks -> score delivered", 8.5, False, MUT)]])
seg = [("gateway",30,PANEL,INK),("ingest",5,PANEL,INK),("jitter",60,GRNF,GRN),
       ("window fill",100,BLUEF,BLUE),("grpc",10,PANEL,INK),("inference",200,VIOF,VIO),
       ("reorder",20,BLUEF,BLUE),("fanout",10,PANEL,INK),("slack",65,AMBF,AMB)]
BY, BH = HB3 + 0.42, 0.66
cx = 0.42
for nm, ms, fl_c, ln_c in seg:
    w = LW_ * ms / 500.0
    sb = box(s3, cx, BY, w, BH, ln_c, fl_c, 1.5, radius=0.02, anchor=MSO_ANCHOR.MIDDLE)
    if w > 0.46: rich(sb, [[(str(ms), 10, True, INK)]], PP_ALIGN.CENTER)
    cx += w
_, lh = free(s3, 0.42, BY + BH + 0.12, LW_, [
 [("jitter 60 ms", 8.5, True, GRN), ("  ·  ", 8.5, False, MUT), ("window fill 100 ms", 8.5, True, BLUE),
  ("  ·  ", 8.5, False, MUT), ("inference 200 ms", 8.5, True, VIO), ("  ·  ", 8.5, False, MUT),
  ("slack 65 ms", 8.5, True, AMB)],
 [("Inference is 46% of the entire budget — so jitter-buffer depth is my single biggest lever before I must start dropping windows.", 8.5, False, INK)],
 [("Every stage is measured separately. One end-to-end number cannot tell you which of the four of us regressed, and each of us will assume it was one of the other three.", 8.5, False, MUT)],
 [("Time-to-first-score is ~2.4 s, not 500 ms — a full 2 s window must exist before anything can be scored.", 8.5, True, RED)]])

opb = [[("OPERATING-POINT TARGETS", 9, True, GRN), ("   reported by name, not as an abstract EER", 8, False, MUT)],
       [("FAR < 1.5%", 8.5, True, INK), ("  synthetic voices that get through", 8, False, MUT)],
       [("FRR < 3.0%", 8.5, True, INK), ("  genuine callers wrongly flagged", 8, False, MUT)],
       [("RTF <= 0.1", 8.5, True, INK), ("  1 s of audio processed in <= 100 ms", 8, False, MUT)],
       [("E2E < 500 ms", 8.5, True, INK), ("  conversational turn-taking budget", 8, False, MUT)],
       [("Fusion model measures FAR 0.55% / FRR 0.55% in-domain — inside target, not yet validated out-of-domain.", 8, False, MUT)]]
OY = BY + BH + 0.12 + lh + 0.18
box(s3, 0.42, OY, LW_, blocks_h(opb, LW_), GRN, GRNP, 2.0, anchor=MSO_ANCHOR.TOP)
rich(s3.shapes[-1], opb)

free(s3, 6.90, HB3 + 0.14, 6.0, [[("DEGRADATION LADDER", 9.5, True, INK),
                            ("   hysteretic: exit only after 10 s of health", 8.5, False, MUT)]])
lad = [("L0  NORMAL","2 s window / 100 ms hop · full cadence",GRN,GRNF),
       ("L1  ELEVATED","queue > 2 windows -> hop widens to 200 ms",GRN,GRNP),
       ("L2  DEGRADED","ring dropping -> hop 500 ms, drop oldest",AMB,AMBF),
       ("L3  DETECTOR DOWN","gRPC gone > 5 s -> scoring stops, call lives",RED,REDF),
       ("L4  SHEDDING","node at capacity -> reject NEW sessions only",RED,REDP)]
ly = HB3 + 0.42
for nm, desc, ln_c, fl_c in lad:
    b = [[(nm + "    ", 9, True, ln_c), (desc, 8.5, False, INK)]]
    h = blocks_h(b, 6.0)
    sh = box(s3, 6.90, ly, 6.0, h, ln_c, fl_c, 2.25, radius=0.14, anchor=MSO_ANCHOR.MIDDLE)
    rich(sh, b); ly += h + 0.12
grb = [[("GOVERNING RULE", 9, True, RED), ("   absence of a score is never equivalent to a low score.", 9, False, INK)],
       [("Reporting “genuine” when the pipeline is broken is the one failure mode that actively causes harm — it green-lights the fraud we exist to stop.", 8, False, MUT)]]
box(s3, 6.90, ly + 0.08, 6.0, blocks_h(grb, 6.0), RED, REDP, 2.25, anchor=MSO_ANCHOR.MIDDLE)
rich(s3.shapes[-1], grb)

SFY = max(OY + blocks_h(opb, LW_), ly + 0.08 + blocks_h(grb, 6.0)) + 0.24
free(s3, 0.42, SFY, 12.5, [[("THE THREE THAT FAIL SILENTLY", 10, True, RED),
  ("   — no exception, no error, just quietly wrong output. 71 edge cases catalogued in total.", 9.5, False, MUT)]])
sf = [
 ("AU-04","AUDIO","Averaging a dual-leg fork",
  "The gateway forks caller and callee as two channels. Averaging them mixes two speakers into one signal, destroying both the deepfake and the speaker-verification signal.",
  "FIX  mono mandatory. Two legs arrive as two sessions sharing a call_id. Never average.",
  "COST IF MISSED  every score on every dual-leg call is meaningless, and nothing reports an error."),
 ("AU-08","AUDIO","Packet-loss concealment",
  "On a bad line the gateway invents audio to cover lost packets. That audio is literally synthetic — so the detector flags a genuine caller as a deepfake.",
  "FIX  gateway sets concealed=true; windows over ~15% concealed samples are skipped or down-weighted.",
  "COST IF MISSED  genuine customers on poor lines get accused of fraud. This is what gets the system switched off."),
 ("SS-06","SESSION","Missing speaker enrolment",
  "No reference voiceprint makes similarity return 0.0, which fusion reads as 'total mismatch' — fabricating fraud signal for every caller who never enrolled.",
  "FIX  emit null plus speaker_signal=UNAVAILABLE. The risk engine renormalises over remaining signals.",
  "COST IF MISSED  every first-time caller looks like an impostor — false accusations at the worst moment."),
]
sf_blk = [[[(c, 9, True, RED), ("   " + k, 8, False, MUT)], [(h, 11, True, INK)],
           [(b, 8, False, INK)], [(f, 8, True, RED)], [(m, 8, True, INK)]]
          for c, k, h, b, f, m in sf]
SFW = 4.06
SFH = max(blocks_h(b, SFW) for b in sf_blk)
sx2 = 0.42
for b in sf_blk:
    sh = box(s3, sx2, SFY + 0.28, SFW, SFH, RED, REDP, 2.25, anchor=MSO_ANCHOR.TOP)
    rich(sh, b); sx2 += 4.22
print("slide 3 ok  bottom=%.2f" % (SFY + 0.28 + SFH))

OUT = "docs/streaming_bridge/VoiceShield-Bridge-3slides.pptx"
prs.save(OUT); print("saved ->", OUT)
